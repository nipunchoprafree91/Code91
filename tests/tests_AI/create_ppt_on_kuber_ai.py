

from pptx import Presentation

# Create a PowerPoint presentation object
presentation = Presentation()

# Slide 1: Title Slide
slide_1 = presentation.slides.add_slide(presentation.slide_layouts[0])
title = slide_1.shapes.title
subtitle = slide_1.placeholders[1]
title.text = "Kubernetes and Docker: Deploying Scalable AI Workloads"
subtitle.text = "An Overview of Containerization and Orchestration for AI Workflows"

# Slide 2: Introduction
slide_2 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide_2.shapes.title
title.text = "Introduction"
content = slide_2.shapes.placeholders[1].text_frame
content.text = "• Kubernetes and Docker are two essential technologies for deploying scalable AI workloads."
p = content.add_paragraph()
p.text = "• Docker helps in containerizing AI applications, while Kubernetes orchestrates and manages containers."
p = content.add_paragraph()
p.text = "• These tools enable flexibility, scalability, and simplified deployment in AI development and production."

# Slide 3: Docker for AI Workloads
slide_3 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide_3.shapes.title
title.text = "Docker for AI Workloads"
content = slide_3.shapes.placeholders[1].text_frame
content.text = "• Docker provides containerization for isolated, portable environments."
p = content.add_paragraph()
p.text = "• Simplifies dependency management and replicability in AI models."
p = content.add_paragraph()
p.text = "• Docker images can include AI frameworks like TensorFlow, PyTorch, and more."
p = content.add_paragraph()
p.text = "• Easily move AI workloads across environments: local, cloud, and edge."

# Slide 4: Kubernetes Overview
slide_4 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide_4.shapes.title
title.text = "Kubernetes Overview"
content = slide_4.shapes.placeholders[1].text_frame
content.text = "• Kubernetes is a container orchestration platform."
p = content.add_paragraph()
p.text = "• Automates deployment, scaling, and operations of containerized applications."
p = content.add_paragraph()
p.text = "• Key components: Pods, Nodes, ReplicaSets, and Services."
p = content.add_paragraph()
p.text = "• Manages distributed AI workloads efficiently."

# Slide 5: Benefits of Kubernetes for AI
slide_5 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide_5.shapes.title
title.text = "Benefits of Kubernetes for AI Workloads"
content = slide_5.shapes.placeholders[1].text_frame
content.text = "• Scalability: Automatically scale up or down AI workloads."
p = content.add_paragraph()
p.text = "• Resource Efficiency: Optimizes resource allocation using autoscaling."
p = content.add_paragraph()
p.text = "• Resilience: Ensures high availability of AI services."
p = content.add_paragraph()
p.text = "• Flexibility: Run AI workloads in hybrid, multi"
